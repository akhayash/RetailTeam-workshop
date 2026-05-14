"""Verify the built solution-architecture.pptx meets acceptance criteria."""
import sys
import warnings
from pathlib import Path

# Capture python-pptx warnings raised while loading
warnings.simplefilter("always")

from pptx import Presentation  # noqa: E402

OUTPUT = Path(
    r"c:\Repos\RetailTeam-workshop\.copilot-tracking\ppt\2026-05-13"
    r"\solution-architecture\slide-deck\solution-architecture.pptx"
)

caught = []
with warnings.catch_warnings(record=True) as w:
    warnings.simplefilter("always")
    prs = Presentation(str(OUTPUT))
    caught.extend(w)

slide_count = len(prs.slides)
width_emu = prs.slide_width
height_emu = prs.slide_height
file_size = OUTPUT.stat().st_size

print(f"file: {OUTPUT}")
print(f"file_size_bytes: {file_size}")
print(f"slide_count: {slide_count}")
print(f"slide_width_emu: {width_emu}")
print(f"slide_height_emu: {height_emu}")
print(f"slide_width_inches: {width_emu / 914400:.3f}")
print(f"slide_height_inches: {height_emu / 914400:.3f}")

expected_w = 12192000
expected_h = 6858000
dims_ok = width_emu == expected_w and height_emu == expected_h
print(f"dimensions_match_13.333x7.5: {dims_ok}")

print("\nspeaker_notes_audit:")
missing = []
empty = []
for idx, slide in enumerate(prs.slides, start=1):
    has_notes_slide = slide.has_notes_slide
    if not has_notes_slide:
        missing.append(idx)
        print(f"  slide {idx:>2}: NO notes_slide")
        continue
    text = slide.notes_slide.notes_text_frame.text.strip()
    char_count = len(text)
    if char_count == 0:
        empty.append(idx)
        print(f"  slide {idx:>2}: notes_slide present but EMPTY")
    else:
        snippet = text[:60].replace("\n", " ")
        print(f"  slide {idx:>2}: {char_count:>4} chars  | {snippet}")

print(f"\nslides_missing_notes_slide: {missing}")
print(f"slides_with_empty_notes: {empty}")
print(f"all_slides_have_nonempty_notes: {not missing and not empty}")

print(f"\npython_pptx_warnings_count: {len(caught)}")
for warning in caught:
    print(f"  WARNING: {warning.category.__name__}: {warning.message}")

ok = (
    slide_count == 15
    and dims_ok
    and not missing
    and not empty
)
print(f"\nACCEPTANCE_PASS: {ok}")
sys.exit(0 if ok else 1)
