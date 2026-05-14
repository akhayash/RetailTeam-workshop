"""Manual validation checks for the solution-architecture deck."""
import json
from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation('../solution-architecture-v2.pptx')

results = {
    "dimensions": {},
    "slide_count": len(prs.slides),
    "speaker_notes": {},
    "element_counts": {},
    "color_samples": {},
}

# 1. Dimensions
w_in = prs.slide_width / 914400
h_in = prs.slide_height / 914400
results["dimensions"] = {
    "width_inches": round(w_in, 3),
    "height_inches": round(h_in, 3),
    "match_expected": abs(w_in - 13.333) < 0.01 and abs(h_in - 7.5) < 0.01,
}
print(f"Slide dimensions: {w_in:.3f}\" x {h_in:.3f}\"")
print(f"Expected: 13.333\" x 7.500\" → {'PASS' if results['dimensions']['match_expected'] else 'FAIL'}")
print(f"Total slides: {len(prs.slides)}")
print()

# 2. Speaker notes
missing_notes = []
for i, slide in enumerate(prs.slides, 1):
    has_notes = False
    if slide.has_notes_slide:
        text = slide.notes_slide.notes_text_frame.text.strip()
        if text:
            has_notes = True
    if not has_notes:
        missing_notes.append(i)

if missing_notes:
    print(f"FAIL: Slides missing speaker notes: {missing_notes}")
else:
    print(f"All {len(prs.slides)} slides have speaker notes ✓")
results["speaker_notes"]["missing"] = missing_notes
print()

# 3. Element counts per slide (focus on diagram slides)
diagram_slides = [5, 6, 13, 14, 15, 16, 17]
print("Element counts (diagram slides):")
for i, slide in enumerate(prs.slides, 1):
    shapes = len(slide.shapes)
    results["element_counts"][i] = shapes
    if i in diagram_slides:
        # Count shape types
        text_boxes = sum(1 for s in slide.shapes if s.has_text_frame)
        auto_shapes = 0
        connectors = 0
        for s in slide.shapes:
            try:
                if hasattr(s, 'auto_shape_type'):
                    _ = s.auto_shape_type
                    auto_shapes += 1
            except (ValueError, AttributeError):
                pass
            if s.shape_type and 'CONNECTOR' in str(s.shape_type):
                connectors += 1
        print(f"  Slide {i:2d}: {shapes:3d} shapes ({text_boxes} text, {auto_shapes} auto-shapes, {connectors} connectors)")

print()

# 4. Check for potential text overflow on dense slides
dense_slides = [6, 14, 16]
print("Dense slide analysis (potential overflow risks):")
for i, slide in enumerate(prs.slides, 1):
    if i not in dense_slides:
        continue
    total_text_len = 0
    small_boxes = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            total_text_len += len(text)
            # Check for very small text boxes with lots of text
            if shape.width and shape.height:
                w_in = shape.width / 914400
                h_in = shape.height / 914400
                area = w_in * h_in
                if area < 2.0 and len(text) > 100:
                    small_boxes += 1
    print(f"  Slide {i:2d}: {total_text_len} chars total, {small_boxes} potentially cramped boxes")

print()

# 5. Color check - sample some fills
print("Color sampling (checking palette compliance):")
target_colors = {
    "14366E": "navy",
    "9BBB59": "olive",
}
found_colors = set()
for i, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if hasattr(shape, 'fill') and shape.fill and shape.fill.type is not None:
            try:
                if shape.fill.fore_color and shape.fill.fore_color.type is not None:
                    rgb = str(shape.fill.fore_color.rgb)
                    found_colors.add(rgb)
                    if rgb.upper() in target_colors:
                        pass  # expected
            except Exception:
                pass

palette_found = []
for hex_val, name in target_colors.items():
    present = hex_val.upper() in {c.upper() for c in found_colors}
    palette_found.append(f"  {name} (#{hex_val}): {'found ✓' if present else 'NOT FOUND'}")
    
for line in palette_found:
    print(line)

results["color_samples"] = {
    "unique_fill_colors_found": len(found_colors),
    "palette_check": palette_found,
}

# Write results
with open("manual-validation-results.json", "w") as f:
    json.dump(results, f, indent=2)

print()
print("Manual validation complete. Results saved to manual-validation-results.json")
