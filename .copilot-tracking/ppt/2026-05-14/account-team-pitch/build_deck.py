"""
Architect-led pitch deck for the Microsoft account team.
Audience: Walmart account team (AE, ATU, STU, ISD sales).
Goal: convince them VirtualMirror AI is a high-quality, sellable engagement
backed by engineering facts.

Run: py build_deck.py
Output: ./slide-deck/virtualmirror-account-team-pitch.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ---------- Theme ----------
NAVY = RGBColor(0x0F, 0x2A, 0x4A)
TEAL = RGBColor(0x00, 0x8B, 0x8B)
ACCENT = RGBColor(0x00, 0x78, 0xD4)
LIGHT = RGBColor(0xF3, 0xF6, 0xFA)
DARK = RGBColor(0x1F, 0x1F, 0x1F)
MUTED = RGBColor(0x55, 0x60, 0x6E)
WARN = RGBColor(0xC2, 0x3B, 0x22)
GOOD = RGBColor(0x10, 0x7C, 0x10)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Segoe UI"
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=DARK, bullet_color=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    bullet_color = bullet_color or TEAL
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        r1 = p.add_run()
        r1.text = "▎ "
        r1.font.size = Pt(size)
        r1.font.color.rgb = bullet_color
        r1.font.bold = True
        r1.font.name = "Segoe UI"
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
        r2.font.name = "Segoe UI"
    return tb


def add_header(slide, eyebrow, title):
    # Top navy band
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.9), NAVY)
    add_text(slide, Inches(0.5), Inches(0.08), Inches(11), Inches(0.35),
             eyebrow, size=11, bold=True, color=RGBColor(0xC2, 0xE0, 0xFF))
    add_text(slide, Inches(0.5), Inches(0.32), Inches(12), Inches(0.55),
             title, size=24, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))


def add_footer(slide, page, total):
    add_rect(slide, 0, Inches(7.18), SLIDE_W, Inches(0.32), NAVY)
    add_text(slide, Inches(0.5), Inches(7.20), Inches(8), Inches(0.28),
             "Microsoft Confidential  |  Architect-led pitch  |  VirtualMirror AI for Walmart",
             size=9, color=RGBColor(0xC2, 0xE0, 0xFF))
    add_text(slide, Inches(11.5), Inches(7.20), Inches(1.5), Inches(0.28),
             f"{page} / {total}", size=9, color=RGBColor(0xC2, 0xE0, 0xFF),
             align=PP_ALIGN.RIGHT)


def add_kpi(slide, x, y, w, h, value, label, color=ACCENT):
    add_rect(slide, x, y, w, h, LIGHT, line=color)
    add_text(slide, x, y + Inches(0.08), w, Inches(0.7), value,
             size=28, bold=True, color=color, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x, y + Inches(0.85), w, Inches(0.45), label,
             size=11, color=MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)


def make_table(slide, x, y, w, h, headers, rows, *, header_fill=NAVY,
               header_color=RGBColor(0xFF, 0xFF, 0xFF), font_size=11):
    cols = len(headers)
    nrows = len(rows) + 1
    tbl_shape = slide.shapes.add_table(nrows, cols, x, y, w, h)
    tbl = tbl_shape.table
    for i, head in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = head
        run.font.bold = True
        run.font.size = Pt(font_size)
        run.font.color.rgb = header_color
        run.font.name = "Segoe UI"
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if ri % 2 == 1 else RGBColor(0xFF, 0xFF, 0xFF)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(font_size)
            run.font.color.rgb = DARK
            run.font.name = "Segoe UI"
    return tbl


# ---------- Build deck ----------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]
TOTAL = 14


def new_slide():
    return prs.slides.add_slide(BLANK)


# Slide 1 — Title
s = new_slide()
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s, 0, Inches(2.4), SLIDE_W, Inches(0.08), TEAL)
add_text(s, Inches(0.7), Inches(0.6), Inches(12), Inches(0.45),
         "MICROSOFT INDUSTRY SOLUTIONS DELIVERY  ·  ARCHITECT BRIEFING",
         size=12, bold=True, color=RGBColor(0xC2, 0xE0, 0xFF))
add_text(s, Inches(0.7), Inches(1.1), Inches(12), Inches(1.4),
         "VirtualMirror AI for Walmart",
         size=44, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(s, Inches(0.7), Inches(2.7), Inches(12), Inches(0.7),
         "Why this is a high-quality engagement to sell — engineering-grade evidence",
         size=22, color=RGBColor(0xC2, 0xE0, 0xFF))
add_text(s, Inches(0.7), Inches(4.0), Inches(12), Inches(2.5),
         ["Audience: Walmart account team (AE, ATU, STU, ISD sales)",
          "Presenter: Solution Architect, Microsoft ISD",
          "Asks: green-light pursuit, executive intro, PTU pre-allocation",
          "Reference: docs/Virtual-Mirror-SOW.md  v0.3.0  (2026-05-14)"],
         size=16, color=RGBColor(0xE6, 0xEE, 0xF7))
add_text(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.4),
         "14 slides  ·  5-minute read  ·  Architect-owned",
         size=11, color=RGBColor(0x9A, 0xB4, 0xD4))

# Slide 2 — The Walmart pain
s = new_slide()
add_header(s, "1. THE PAIN", "Walmart's apparel returns problem is structural, not seasonal")
add_kpi(s, Inches(0.5), Inches(1.2), Inches(2.9), Inches(1.3), "$14.7B", "Walmart online apparel revenue (2024)")
add_kpi(s, Inches(3.6), Inches(1.2), Inches(2.9), Inches(1.3), "24–26%", "Industry online apparel return rate")
add_kpi(s, Inches(6.7), Inches(1.2), Inches(2.9), Inches(1.3), "53–70%", "Of returns driven by FIT, not look", color=WARN)
add_kpi(s, Inches(9.8), Inches(1.2), Inches(2.9), Inches(1.3), "$200–400M", "Annual avoidable cost (fit-only)", color=WARN)

add_text(s, Inches(0.5), Inches(2.8), Inches(12.3), Inches(0.45),
         "What's broken today", size=16, bold=True, color=NAVY)
add_bullets(s, Inches(0.5), Inches(3.25), Inches(12.3), Inches(2.0), [
    "Zeekit answers \"how does it look on me?\" — it does NOT answer \"will it fit my body?\"",
    "Shoppers guess sizes. They buy 2–3 sizes and return the rest. Reverse logistics + restocking + markdowns + write-offs.",
    "Static size charts vary by brand and supplier. There is no measurement-based confidence layer at the PDP.",
    "Returns erode apparel margin twice: once in cost-to-return, once in markdown of returned inventory.",
], size=14)

add_rect(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.4), LIGHT, line=ACCENT)
add_text(s, Inches(0.7), Inches(5.6), Inches(12), Inches(0.4),
         "Architect's read", size=13, bold=True, color=ACCENT)
add_text(s, Inches(0.7), Inches(5.95), Inches(12), Inches(0.9),
         "This is a measurable, instrumentable problem with a clean technology wedge. "
         "The pain is large enough to fund a multi-year platform play, and the gap (measurement vs. visualization) "
         "is genuinely unaddressed by Walmart's current Zeekit investment.",
         size=13, color=DARK)
add_footer(s, 2, TOTAL)

# Slide 3 — The opportunity
s = new_slide()
add_header(s, "2. THE OPPORTUNITY", "Even modest reductions translate to nine-figure outcomes")

make_table(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(2.6),
           ["Driver", "Basis", "Conservative", "Target"],
           [
               ["Fit-driven return cost reclaim", "$200–400M addressable",        "10% reduction = $20–40M",   "30% reduction = $60–120M"],
               ["Fit-driven returns (industry math)", "$14.7B × 25–40% × ~70%",   "$2.6B–$4.1B surface",       "30% cut = $0.8B–$1.2B savings"],
               ["Conversion lift on enabled SKUs",  "+3–5% (BO-2 target)",        "+3% on enabled assortment", "+5% on enabled assortment"],
               ["Market-share reclaim potential",   "+300–500 bps vs. baseline",  "~$440M incremental rev",    "~$735M incremental rev"],
               ["Cost per assessment",              "Azure consumption model",     "$0.018 (pilot)",            "$0.012 (growth)"],
           ], font_size=11)

add_rect(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(1.6), LIGHT, line=ACCENT)
add_text(s, Inches(0.7), Inches(4.1), Inches(12), Inches(0.4),
         "Unit economics tell the story", size=13, bold=True, color=ACCENT)
add_text(s, Inches(0.7), Inches(4.45), Inches(12), Inches(1.1),
         "At $0.012–$0.018 per assessment and ~$48K Year-1 Azure consumption, ROI breakeven sits "
         "well under 0.1% of fit-related avoidable cost. The economics are not the deal blocker — "
         "Walmart's adoption velocity is. That is what this engagement is engineered to unlock.",
         size=13, color=DARK)

add_text(s, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.4),
         "Source: docs/Virtual-Mirror-SOW.md §2.1, §3.1; cost-estimate.md; Product-definition.md — Market Sizing.",
         size=10, color=MUTED)
add_footer(s, 3, TOTAL)

# Slide 4 — Why Walmart, why now
s = new_slide()
add_header(s, "3. WHY WALMART, WHY NOW", "A clean technology wedge against a stated executive priority")

add_text(s, Inches(0.5), Inches(1.15), Inches(6.0), Inches(0.4),
         "Walmart-specific signals", size=15, bold=True, color=NAVY)
add_bullets(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(4.5), [
    "3rd-largest U.S. apparel e-com retailer — scale matches our delivery shape.",
    "Already invested in Zeekit (visual try-on). VirtualMirror is the complementary measurement layer, not a competitive overlap.",
    "Apparel returns are publicly named as a margin pressure in Walmart Digital's investor narrative.",
    "Walmart Platform Engineering operates a golden-path Azure landing zone — our Bicep + ACA pattern lands cleanly.",
    "Privacy Office is mature; DPIA cycles are repeatable. Lowers our delivery risk.",
], size=13)

add_text(s, Inches(7.0), Inches(1.15), Inches(5.8), Inches(0.4),
         "Why now (12-month window)", size=15, bold=True, color=NAVY)
add_bullets(s, Inches(7.0), Inches(1.55), Inches(5.8), Inches(4.5), [
    "GPT-5.2 Vision GA in East US 2 enables single-photo measurement at production cost.",
    "Azure AI Content Safety + Foundry Florence-2 close the responsible-AI surface (minor detection, content moderation).",
    "EU AI Act enforcement ramp creates a compliance moat for retailers who design responsibly NOW.",
    "Competitor virtual-fit startups (3DLOOK, Bold Metrics) are point solutions; Microsoft offers an Azure-native, multi-tenant platform.",
    "First-mover signal: a successful v1 becomes a Microsoft-Walmart Industry reference, unlocking land-and-expand into Sam's Club and marketplace sellers.",
], size=13)

add_footer(s, 4, TOTAL)

# Slide 5 — What we're selling
s = new_slide()
add_header(s, "4. WHAT WE ARE SELLING", "An API-first, multi-tenant clothing fit assessment service on Azure")

add_rect(s, Inches(0.5), Inches(1.15), Inches(4.0), Inches(2.6), LIGHT, line=ACCENT)
add_text(s, Inches(0.65), Inches(1.25), Inches(3.7), Inches(0.4),
         "Core capability", size=13, bold=True, color=ACCENT)
add_bullets(s, Inches(0.65), Inches(1.65), Inches(3.7), Inches(2.0), [
    "POST /v1/assessments → 5-point fit per area + confidence",
    "POST /v1/profiles → opt-in, 24h hard-delete",
    "POST /v1/garments → single + batch (≤100), version history",
    "OAuth 2.0 / Entra ID, per-operation scopes",
], size=12)

add_rect(s, Inches(4.65), Inches(1.15), Inches(4.0), Inches(2.6), LIGHT, line=TEAL)
add_text(s, Inches(4.8), Inches(1.25), Inches(3.7), Inches(0.4),
         "Three-tier AI pipeline", size=13, bold=True, color=TEAL)
add_bullets(s, Inches(4.8), Inches(1.65), Inches(3.7), Inches(2.0), [
    "Tier 1: Florence-2 (people detection, bounding box)",
    "Tier 1: Content Safety (minor + inappropriate)",
    "Tier 2: GPT-5.2 Vision (structured measurements)",
    "Tier 3 (v2 upsell): SMPL custom model on Foundry",
], size=12, bullet_color=TEAL)

add_rect(s, Inches(8.8), Inches(1.15), Inches(4.0), Inches(2.6), LIGHT, line=NAVY)
add_text(s, Inches(8.95), Inches(1.25), Inches(3.7), Inches(0.4),
         "Production hardening", size=13, bold=True, color=NAVY)
add_bullets(s, Inches(8.95), Inches(1.65), Inches(3.7), Inches(2.0), [
    "12 Bicep modules, dev/staging/prod IaC",
    "GitHub Actions: SAST/SCA/SBOM/Trivy/Notation",
    "OTel → Azure Monitor + 3 runbooks",
    "NBomber 500-concurrent + chaos validation",
], size=12, bullet_color=NAVY)

add_rect(s, Inches(0.5), Inches(3.95), Inches(12.3), Inches(2.9), LIGHT, line=MUTED)
add_text(s, Inches(0.65), Inches(4.05), Inches(12), Inches(0.4),
         "Architecture (one-line)", size=13, bold=True, color=NAVY)
add_text(s, Inches(0.65), Inches(4.45), Inches(12), Inches(0.6),
         "Frontend → Front Door/WAF → ACA (.NET 8 API) → [Florence-2 · Content Safety · GPT-5.2 Vision] → Cosmos DB · Blob (60s TTL) · Service Bus · Key Vault → Azure Monitor",
         size=12, color=DARK)
add_text(s, Inches(0.65), Inches(5.15), Inches(12), Inches(0.4),
         "What is intentionally OUT of v1 (upsell map)", size=13, bold=True, color=NAVY)
add_bullets(s, Inches(0.65), Inches(5.55), Inches(12), Inches(1.3), [
    "Multi-region active-active  ·  Custom SMPL Tier-3 model  ·  Native mobile SDK / on-device inference",
    "Microsoft Fabric intelligence loop (Concept C)  ·  Direct Zeekit integration  ·  Supplier-feed ingestion",
    "Storefront PDP UI / A-B harness  ·  Production support beyond 30-day hypercare",
], size=12)

add_footer(s, 5, TOTAL)

# Slide 6 — Engineering proof points
s = new_slide()
add_header(s, "5. ENGINEERING PROOF POINTS", "Hypothesis-driven delivery — every claim has a measurable gate")

make_table(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(3.5),
           ["Hyp.", "Claim under test", "Gate (acceptance)", "Sprint", "If it fails"],
           [
               ["H1", "GPT-5.2 Vision measurement accuracy ±2–4 cm", "≤15% of ≥50-image diverse calibration set > ±4 cm", "End S3", "3DLOOK bridge (+2 wks, Change Order) — known fallback"],
               ["H7", "AI failover < 5 s under primary outage",       "Polly chaos test trips circuit, fallback responds < 5s", "End S2", "Re-tune Polly thresholds; degrade to size-chart L5"],
               ["H3", "p95 latency < 5 s end-to-end",                  "NBomber sustained 30-min run @ 500 concurrent", "End S8", "Add KEDA scale-out, raise OpenAI PTU"],
               ["H5", "500-concurrent capacity",                       "KEDA scales 2→8 instances, no OOM",               "End S8", "Provision additional ACA replicas; PTU bump"],
               ["H8", "≥90% success during chaos",                     "Multi-fault inject (AI+Cosmos+SB) over 10 min",   "End S8", "Strengthen circuit breakers + DLQ replay runbook"],
           ], font_size=11)

add_rect(s, Inches(0.5), Inches(4.85), Inches(12.3), Inches(1.95), LIGHT, line=GOOD)
add_text(s, Inches(0.7), Inches(4.95), Inches(12), Inches(0.4),
         "Why this matters to Walmart (and to you in the sales cycle)", size=13, bold=True, color=GOOD)
add_bullets(s, Inches(0.7), Inches(5.35), Inches(12), Inches(1.4), [
    "Each hypothesis has a date, a metric, and a defined contingency. There are no \"trust us\" clauses in the SOW.",
    "MVP gate at end of Sprint 4 (8 weeks). Walmart sees a working API, not a slideware demo.",
    "Hypothesis failures invoke pre-priced contingencies — predictable cost ceiling for the Customer.",
    "Engineering rigor is the differentiator. Competitors will pitch demos; we pitch acceptance criteria.",
], size=12)
add_footer(s, 6, TOTAL)

# Slide 7 — Risk-adjusted delivery
s = new_slide()
add_header(s, "6. RISK-ADJUSTED DELIVERY", "Top 7 risks have named owners, mitigations, and validation steps")

make_table(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(4.6),
           ["ID", "Risk", "Sev", "Mitigation", "How we prove it's controlled"],
           [
               ["PR-1", "H1 measurement accuracy fails",          "HIGH",     "Pre-S3 prompt spike + calibration set; 3DLOOK bridge", "S3 gate report w/ confidence intervals"],
               ["PR-2", "Walmart Azure tenant provisioning slips","HIGH",     "MSDN fallback for Sprint 1; promote in Sprint 2",       "C-1 received by end S2 (RAG green)"],
               ["PR-3", "Tech Lead PR-review bottleneck",         "HIGH",     "25% reserved review time; cross-review by SDEs",        "PR-cycle-time < 24 h in retro metrics"],
               ["PR-4", "DPIA finds late blockers",               "MED",      "S4 mid-checkpoint + S8 final review",                   "No Open/High findings remain at GA"],
               ["PR-5", "Bicep complexity > estimate",            "MED",      "Full DevOps FTE from S6; Security pair-review",          "what-if clean across 3 envs by S7"],
               ["PR-6", "Catalog feed delayed > S5",              "MED",      "Synthetic catalog covers US4 dev",                       "Fallback trigger if not received by end S6"],
               ["PR-7", ".NET 8 talent gap post-handover",        "HIGH",     "Self-contained API, runbooks, KT (D-14); MS-managed option",  "S6 Council confirms ownership model + ADR"],
           ], font_size=10)

add_rect(s, Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.95), LIGHT, line=NAVY)
add_text(s, Inches(0.7), Inches(6.05), Inches(12), Inches(0.4),
         "Account-team takeaway", size=13, bold=True, color=NAVY)
add_text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.5),
         "PR-7 (.NET talent gap) is your post-sale conversation starter for a managed-services SOW. "
         "The risk is named in §13 of the SOW — surface it early, frame it as a planned land-and-expand, not a surprise.",
         size=12, color=DARK)
add_footer(s, 7, TOTAL)

# Slide 8 — Why Microsoft (differentiation)
s = new_slide()
add_header(s, "7. WHY MICROSOFT", "Five differentiators no point-solution vendor can match")

items = [
    ("Azure-native AI stack",
     "GPT-5.2 Vision + Florence-2 + Content Safety from a single trust boundary. "
     "Managed identity end-to-end. No third-party data-processing agreements to negotiate."),
    ("Multi-tenant by design, not afterthought",
     "Cosmos hierarchical partition keys, per-tenant scopes, per-tenant tolerance bands. "
     "Walmart at GA, architected for Sam's Club and marketplace sellers as the v1.x roadmap."),
    ("Responsible-AI as a delivered artifact",
     "Model card, DPIA input package, minor-detection refusal, bias considerations — included in the engagement, "
     "not a separate consulting line."),
    ("Engineering rigor codified",
     "OWASP ASVS L2, SOC 2 mapped, NIST CSF 2.0, OpenTelemetry observability, NBomber + chaos, "
     "SBOM + Trivy + Notation. This is the floor, not the ceiling."),
    ("Land-and-expand path is pre-built",
     "v2 SMPL custom model · multi-region · mobile SDK · Microsoft Fabric intelligence loop. "
     "Each is a sized scope-trap (ST-1..ST-10) in the SOW — pre-priced upsells."),
]
y = Inches(1.15)
for title, body in items:
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.05), LIGHT, line=ACCENT)
    add_text(s, Inches(0.7), y + Inches(0.08), Inches(11.9), Inches(0.4),
             title, size=14, bold=True, color=ACCENT)
    add_text(s, Inches(0.7), y + Inches(0.45), Inches(11.9), Inches(0.55),
             body, size=11, color=DARK)
    y += Inches(1.15)

add_footer(s, 8, TOTAL)

# Slide 9 — Engagement shape
s = new_slide()
add_header(s, "8. ENGAGEMENT SHAPE", "9 sprints · 18 weeks · ~274 PD · fixed capacity, variable scope")

add_kpi(s, Inches(0.5),  Inches(1.15), Inches(2.4), Inches(1.3), "18 wks", "9 × 2-week sprints + 2-wk hypercare")
add_kpi(s, Inches(3.0),  Inches(1.15), Inches(2.4), Inches(1.3), "~274 PD", "Build envelope (T&M cap)", color=TEAL)
add_kpi(s, Inches(5.5),  Inches(1.15), Inches(2.4), Inches(1.3), "5.75 FTE", "Average team (peak 6.75 in S7–8)")
add_kpi(s, Inches(8.0),  Inches(1.15), Inches(2.4), Inches(1.3), "147 tasks", "Pre-decomposed in tasks.md", color=TEAL)
add_kpi(s, Inches(10.5), Inches(1.15), Inches(2.3), Inches(1.3), "$0.012–0.018", "Per-assessment Azure cost", color=GOOD)

make_table(s, Inches(0.5), Inches(2.7), Inches(12.3), Inches(3.6),
           ["Sprint", "Headline output", "Gate"],
           [
               ["S1",   "Solution scaffold, domain + interfaces, auth middleware skeleton",                     "—"],
               ["S2",   "Foundational platform, resilience pipelines, Aspire wired",                            "H7 — AI failover < 5 s"],
               ["S3",   "AI clients online, image validation, accuracy spike result",                           "H1 — measurement ±2–4 cm"],
               ["S4",   "MVP demo: POST /v1/assessments returns 5-point fit",                                   "MVP demo to Walmart"],
               ["S5",   "Profile CRUD + by-profile endpoint; OpenAPI contract validated",                        "—"],
               ["S6",   "Garment ingestion (single + batch); Bicep ~50%",                                        "—"],
               ["S7",   "IaC complete; CI green; alerts + 3 runbooks",                                           "—"],
               ["S8",   "NBomber 500-concurrent, chaos passes, model card, DR plan",                             "H3, H5, H8"],
               ["S9",   "UAT, security + privacy sign-off, canary 5→25→100%",                                    "Production go-live"],
           ], font_size=11)

add_text(s, Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.5),
         "Pricing model: T&M within fixed-capacity envelope. Not-to-exceed clause negotiated at Work Order signature. "
         "Capacity-based agile — Microsoft does not guarantee whole-backlog delivery; Microsoft guarantees the team and capacity.",
         size=11, color=MUTED)

add_footer(s, 9, TOTAL)

# Slide 10 — Commercial envelope
s = new_slide()
add_header(s, "9. COMMERCIAL ENVELOPE", "Clean, defensible, MSA-governed — no hidden tail")

add_text(s, Inches(0.5), Inches(1.15), Inches(6.0), Inches(0.4),
         "Deal structure", size=15, bold=True, color=NAVY)
make_table(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(3.0),
           ["Item", "Value"],
           [
               ["Engagement type",   "Agile, fixed-capacity + variable scope"],
               ["Pricing model",     "T&M w/ NTE (negotiated at Work Order)"],
               ["Validity",          "30 days from 2026-05-14"],
               ["Master agreement",  "Walmart-Microsoft MSA"],
               ["Azure consumption", "Customer-paid, ~$48K Year 1"],
               ["Hypercare",         "30 days post-GA (best-effort, business hours)"],
               ["Completion",        "Capacity OR term OR backlog OR termination"],
           ], font_size=11)

add_text(s, Inches(7.0), Inches(1.15), Inches(5.8), Inches(0.4),
         "What protects Microsoft (and the Customer)", size=15, bold=True, color=NAVY)
add_bullets(s, Inches(7.0), Inches(1.55), Inches(5.8), Inches(5.0), [
    "BO-1/BO-2 are Customer-owned outcomes — Microsoft does not guarantee revenue/return-rate uplift.",
    "Capacity caveat: backlog completeness is NOT guaranteed; capacity IS.",
    "Deemed-acceptance after 5 business days on D-4, D-9, D-13, D-14 prevents review-stall.",
    "Default-no-change rule on Change Orders (3 BD lapse) prevents scope creep through silence.",
    "Non-GA product disclosure (A-7) — GPT-5.2 GA risk transferred contractually.",
    "PR-7 (.NET talent gap) named upfront — sets up managed-services SOW conversation.",
    "Responsible AI sensitive-use review (A-24) — protects MS reputational risk.",
], size=12)

add_footer(s, 10, TOTAL)

# Slide 11 — Land and expand
s = new_slide()
add_header(s, "10. LAND AND EXPAND", "v1 is the wedge — the SOW already itemizes the upsell map")

make_table(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(4.5),
           ["Trap ID", "Upsell opportunity", "Sized estimate", "Trigger"],
           [
               ["ST-1",  "H1 contingency: 3DLOOK bridge integration",            "+20–30 PD, +2 wks",       "If H1 gate misses ±4 cm"],
               ["ST-2",  "Multi-region active-active (v2)",                       "+60–90 PD",               "Walmart raises SLA target to 99.95%+"],
               ["ST-3",  "Native mobile SDK + on-device inference (Concept B)",   "+120 PD+",                "In-store / fitting-room expansion"],
               ["ST-4",  "Each additional tenant beyond Walmart",                 "+5–10 PD per tenant",     "Sam's Club, marketplace sellers"],
               ["ST-5",  "Catalog data quality remediation",                      "+10–25 PD",               "Supplier-feed normalization"],
               ["ST-6",  "Storefront PDP widget / SDK reference impl",            "+15–25 PD",               "Walmart frontend org adopts directly"],
               ["ST-7",  "HIPAA / FedRAMP / EU AI Act high-risk compliance",      "+40 PD+",                 "Regulatory expansion or marketplace trust"],
               ["ST-8",  "Walmart-storefront SLO performance acceptance",         "+10–15 PD",               "Beyond NBomber smoke"],
               ["ST-9",  "Production support / on-call beyond hypercare",         "Separate managed-svcs SOW","Default outcome — not optional"],
               ["ST-10", "Direct Zeekit integration",                             "+20–40 PD",               "Partnership stand-up"],
           ], font_size=10)

add_rect(s, Inches(0.5), Inches(5.85), Inches(12.3), Inches(1.05), LIGHT, line=GOOD)
add_text(s, Inches(0.7), Inches(5.95), Inches(12), Inches(0.4),
         "Land-and-expand math", size=13, bold=True, color=GOOD)
add_text(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.6),
         "Initial v1 + ST-9 (managed services) + 1–2 likely traps (ST-1, ST-3 or ST-6) doubles deal size within 12 months. "
         "This is engineered into the deal shape, not bolted on.",
         size=12, color=DARK)
add_footer(s, 11, TOTAL)

# Slide 12 — What we need from the account team
s = new_slide()
add_header(s, "11. WHAT WE NEED FROM YOU", "Five concrete asks to close the pursuit")

asks = [
    ("Executive sponsor introduction",
     "Walmart Digital VP (Apparel) or SVP Customer Experience. Architect + Engagement Lead in the room."),
    ("Privacy Office pre-engagement",
     "30-minute discovery call to confirm DPIA cadence and minor-detection posture. De-risks PR-4 before signature."),
    ("Azure subscription + Entra ID confirmation",
     "Confirm landing-zone availability and PTU pre-allocation discussion with the Azure account team. De-risks PR-2 and C-1/C-2."),
    ("Calibration dataset feasibility",
     "Confirm whether Walmart can provide ≥100 photos with ground-truth measurements (anonymized internal OK). "
     "If not, we pre-plan a Microsoft-procured set in Sprint 1. De-risks H1."),
    ("Pursuit support: 30 days to countersignature",
     "SOW v0.3.0 validity expires 2026-06-13. Need account-team commitment to drive procurement and legal review on that timeline."),
]
y = Inches(1.15)
for i, (title, body) in enumerate(asks, start=1):
    add_rect(s, Inches(0.5), y, Inches(0.6), Inches(1.05), ACCENT)
    add_text(s, Inches(0.5), y, Inches(0.6), Inches(1.05), str(i),
             size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(1.1), y, Inches(11.7), Inches(1.05), LIGHT, line=ACCENT)
    add_text(s, Inches(1.3), y + Inches(0.08), Inches(11.4), Inches(0.4),
             title, size=14, bold=True, color=NAVY)
    add_text(s, Inches(1.3), y + Inches(0.45), Inches(11.4), Inches(0.55),
             body, size=11, color=DARK)
    y += Inches(1.13)

add_footer(s, 12, TOTAL)

# Slide 13 — Why this is sellable (architect's verdict)
s = new_slide()
add_header(s, "12. ARCHITECT'S VERDICT", "Why I am confident this is a sellable, deliverable engagement")

add_text(s, Inches(0.5), Inches(1.15), Inches(6.0), Inches(0.4),
         "Sellability — to Walmart", size=15, bold=True, color=GOOD)
add_bullets(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(5.0), [
    "Named pain ($200–400M) + named gap (Zeekit covers look, not fit).",
    "Unit economics ($0.012–0.018/assessment) are non-controversial.",
    "Hypothesis-gated delivery makes the ROI defensible to the Walmart Digital exec team.",
    "Responsible-AI artifacts answer Privacy Office's first three questions.",
    "Multi-tenant architecture aligns to Sam's Club / marketplace strategy.",
], size=13, bullet_color=GOOD)

add_text(s, Inches(7.0), Inches(1.15), Inches(5.8), Inches(0.4),
         "Deliverability — by Microsoft ISD", size=15, bold=True, color=NAVY)
add_bullets(s, Inches(7.0), Inches(1.55), Inches(5.8), Inches(5.0), [
    "147 tasks pre-decomposed; 8 work packages; ~14% headroom.",
    "Senior IC team with named bottleneck mitigation (PR-3).",
    "5 measurable hypothesis gates with pre-priced contingencies.",
    "Bicep + GitHub Actions + OTel — golden-path aligned to Walmart's platform.",
    "MSA-governed, validity-bounded, deemed-acceptance protected.",
    "PR-7 surfaces talent-gap risk early — sets up managed-services SOW.",
], size=13)

add_rect(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.55), GOOD)
add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.55),
         "Recommendation: PURSUE. Engineering risk is bounded. Commercial structure is clean. Land-and-expand path is pre-built.",
         size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 13, TOTAL)

# Slide 14 — Call to action / next steps
s = new_slide()
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s, 0, Inches(2.4), SLIDE_W, Inches(0.08), TEAL)
add_text(s, Inches(0.7), Inches(0.6), Inches(12), Inches(0.45),
         "NEXT 14 DAYS",
         size=12, bold=True, color=RGBColor(0xC2, 0xE0, 0xFF))
add_text(s, Inches(0.7), Inches(1.05), Inches(12), Inches(1.0),
         "Let's get this in front of Walmart.",
         size=36, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(s, Inches(0.7), Inches(2.7), Inches(12), Inches(0.5),
         "Three actions, this week:",
         size=18, bold=True, color=RGBColor(0xC2, 0xE0, 0xFF))
add_bullets(s, Inches(0.7), Inches(3.3), Inches(12), Inches(2.5), [
    "AE: schedule 60-min exec discovery with Walmart Digital VP (Apparel).",
    "AE + ATU: confirm Azure account team + PTU pre-allocation pathway.",
    "Architect (me): walk Walmart Privacy Office through DPIA cadence + minor-detection design.",
], size=16, color=RGBColor(0xE6, 0xEE, 0xF7), bullet_color=TEAL)

add_rect(s, Inches(0.7), Inches(5.7), Inches(12), Inches(1.1), RGBColor(0x14, 0x39, 0x5E), line=TEAL)
add_text(s, Inches(0.9), Inches(5.78), Inches(11.6), Inches(0.4),
         "SOW reference  ·  docs/Virtual-Mirror-SOW.md v0.3.0",
         size=13, bold=True, color=RGBColor(0xC2, 0xE0, 0xFF))
add_text(s, Inches(0.9), Inches(6.15), Inches(11.6), Inches(0.6),
         "147 tasks · 274 PD · 18 weeks · 9 sprints · ~$48K Year-1 Azure · 5 hypothesis gates · 30-day validity expires 2026-06-13.",
         size=12, color=RGBColor(0xE6, 0xEE, 0xF7))

add_text(s, Inches(0.7), Inches(7.0), Inches(12), Inches(0.4),
         "Microsoft Confidential  ·  Architect-led pitch  ·  14 / 14",
         size=10, color=RGBColor(0x9A, 0xB4, 0xD4))


# ---------- Save ----------
OUT = Path(__file__).parent / "slide-deck"
OUT.mkdir(parents=True, exist_ok=True)
out_path = OUT / "virtualmirror-account-team-pitch.pptx"
prs.save(str(out_path))
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
